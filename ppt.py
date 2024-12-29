from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(slide_data, theme='default'):
    prs = Presentation()

    for slide_info in slide_data:
        slide_layout = prs.slide_layouts[1]  # Use the title and content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_info['title']
        content.text = slide_info['content']

        if theme == 'theme1':
            apply_theme1(slide)
        elif theme == 'theme2':
            apply_theme2(slide)
        else:
            apply_default_theme(slide)

    prs.save('generated_presentation.pptx')

def apply_default_theme(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True

def apply_theme1(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.italic = True

def apply_theme2(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.underline = True

if __name__ == "__main__":
    slide_data = [
        {'title': 'Slide 1', 'content': 'Content for slide 1'},
        {'title': 'Slide 2', 'content': 'Content for slide 2'},
        {'title': 'Slide 3', 'content': 'Content for slide 3'}
    ]
    create_presentation(slide_data, theme='theme1')